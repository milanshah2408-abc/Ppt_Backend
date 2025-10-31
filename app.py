from urllib import response
from flask import Flask, request, jsonify, send_file
from flask_cors import CORS
from pptx import Presentation
from pptx.util import Inches, Pt
from asgiref.wsgi import WsgiToAsgi
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
import tempfile
from dotenv import load_dotenv
import json
from uvicorn import run
import requests
import google.generativeai as genai


load_dotenv()

app = Flask(__name__)
CORS(app)
asgi_app = WsgiToAsgi(app)

# Configure Gemini
GOOGLE_API_KEY = os.getenv('GOOGLE_API_KEY')
genai.configure(api_key=GOOGLE_API_KEY)
model = genai.GenerativeModel('gemini-2.5-flash')


@app.route('/api/generate-ppt-from-text', methods=['POST'])
def generate_ppt_from_text():
    try:
        data = request.get_json()
        user_text = data.get('user_text')
        slides_count = int(data.get('slides', 8) or 8)
        theme_file = data.get('theme', 'theme.pptx')
        if not user_text or not user_text.strip():
            return jsonify({'error': 'Custom text is required'}), 400

        # Use Gemini to extract super title and structure slides
        prompt = f"""
        Analyze the following text and generate a super title and a structured PowerPoint outline:
        ---
        {user_text}
        ---
        Requirements:
        - Extract a concise, descriptive super title summarizing the main idea.
        - Produce EXACTLY {slides_count} content slides (title slide is added separately by the app).
        - Each slide should have a clear title and 3-6 informative bullet points.
        - Optionally include speaker notes per slide (1-2 sentences).
        - Return ONLY valid JSON in this schema:
        {{
          "super_title": "string",
          "slides": [
            {{
              "title": "string",
              "content": ["bullet 1", "bullet 2", "bullet 3"],
              "speaker_notes": "optional short paragraph"
            }}
          ]
        }}
        """

        response = model.generate_content(
            prompt,
            generation_config=genai.types.GenerationConfig(
                temperature=0.4,
                max_output_tokens=1400
            )
        )
        raw = response.text
        parsed = try_parse_json(raw)
        if parsed is None:
            raise ValueError("Model returned non-JSON")

        super_title = parsed.get('super_title', 'Presentation')
        slides = parsed.get('slides', [])
        # Normalize to expected format for create_presentation
        normalized = {
            'title': super_title,
            'slides': slides,
            'references': []
        }
        pptx_path = create_presentation(super_title, normalized, theme_file)
        return send_file(
            pptx_path,
            as_attachment=True,
            download_name=f"{super_title.replace(' ', '_')}_presentation.pptx",
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )
    except Exception as e:
        import traceback
        print('Error in /api/generate-ppt-from-text:', traceback.format_exc())
        return jsonify({'error': str(e), 'trace': traceback.format_exc()}), 500

# Endpoint to list available themes
@app.route('/api/themes', methods=['GET'])
def list_themes():
    themes_dir = os.path.join(os.path.dirname(__file__), 'themes')
    files = []
    for f in os.listdir(themes_dir):
        if f.lower().endswith(('.pptx', '.thmx')):
            files.append(f)
    return jsonify({'themes': files})

@app.route('/api/generate-ppt', methods=['POST'])
def generate_ppt():
    try:
        data = request.get_json()
        topic = data.get('topic')
        requested_slides = int(data.get('slides', 8) or 8)
        # Enforce at least 4 content slides (Intro, Advantages, Disadvantages, Conclusion) and max 25
        slides_count = max(4, min(25, requested_slides))
        research_boost = bool(data.get('research', True))
        
        if not topic:
            return jsonify({'error': 'Topic is required'}), 400
        
        # Optional lightweight research (Wikipedia) to ground the content
        research = None
        if research_boost:
            research = fetch_wikipedia_research(topic)
        
        # Generate content using Gemini, guided by research if available
        content = generate_content_with_openai(topic, slides_count, research)
        
    # Get theme from request, default to 'theme.pptx'
        theme_file = data.get('theme', 'theme.pptx')
    # Create PowerPoint presentation
        pptx_path = create_presentation(topic, content, theme_file)
        
        # Return the file
        return send_file(
            pptx_path,
            as_attachment=True,
            download_name=f"{topic.replace(' ', '_')}_presentation.pptx",
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )
    
    except Exception as e:
        return jsonify({'error': str(e)}), 500

def generate_content_with_openai(topic, slides_count, research):


    response_schema = {
        "type": "object",
        "properties": {
            "title": {"type": "string"},
            "slides": {
                "type": "array",
                "items": {
                    "type": "object",
                    "properties": {
                        "title": {"type": "string"},
                        "content": {
                            "type": "array",
                            "items": {"type": "string"}
                        },
                        "speaker_notes": {"type": "string"}
                    },
                    "required": ["title", "content"]
                }
            },
            "references": {
                "type": "array",
                "items": {"type": "string"}
            }
        },
        "required": ["title", "slides"]
    }

    """Generate presentation content using Gemini, grounded with research when available.
    Enforce: Introduction, Advantages, Disadvantages, Conclusion. If more slides requested, add informative sections.
    """
    try:
        research_text = ''
        references = []
        if research and research.get('summary'):
            research_text = research['summary'][:6000]
            if research.get('url'):
                references.append(research['url'])

        prompt = f"""
        You are a professional presentation creator. Build a factual, well-structured PowerPoint outline
        about the topic "{topic}".

        Ground your content in the following research (if present). Do not fabricate facts or statistics.
        Research:
        ---
        {research_text}
        ---

        Requirements:
        - Produce EXACTLY {slides_count} content slides (title slide is added separately by the app).
        - Enforce order and presence:
          1) Introduction (3-5 concise sentences tailored to the topic, explain what it is and why it matters)
          2) Advantages (4-6 specific, fact-grounded benefits with clear explanations)
          3) Disadvantages (4-6 specific, fact-grounded limitations or challenges with clear explanations)
          4) Conclusion (3-4 concise sentences summarizing key takeaways and next steps)
        - If more than 4 slides requested, add additional informative sections BEFORE Conclusion with clear titles like "Key Concepts", "Use Cases", "Examples", "Statistics", "Trends", "Current Status", "Future Outlook", "Implementation", "Best Practices", or "Challenges". Ensure they are relevant to the topic and grounded in the research.
        - Each non-intro/conclusion slide should have 4-6 concise, informative bullets.
        - Make content specific to the topic, not generic. Use research data when available.
        - Include optional speaker_notes per slide (1-3 sentences) to guide the presenter.
        - Include a references array with source URLs at the end (use research URL if available).
        """
        
        print(f"Generating content with Gemini for topic: {topic}, slides: {slides_count}")

        # Use the Gemini API to generate content in JSON format using the defined schema.
        response = genai.GenerativeModel('gemini-1.5-flash').generate_content(
            prompt,
            generation_config=genai.types.GenerationConfig(
                temperature=0.4,
                max_output_tokens=1400,
                response_mime_type="application/json",
                response_schema=response_schema
            )
        )
        
        # Access the text of the response, which is guaranteed to be valid JSON.
        raw = response.text
        print(f"Gemini response received, length: {len(raw)}")
        
        # Directly parse the raw text as JSON. The try_parse_json function is no longer needed
        # because the API is enforcing the format.
        parsed = json.loads(raw)

        # Ensure references include research url if missing
        if references:
            parsed.setdefault('references', [])
            if research['url'] not in parsed['references']:
                parsed['references'].append(research['url'])

        print(f"Successfully generated {len(parsed.get('slides', []))} slides with Gemini")
        
        # Normalize structure strictly to requested format and count
        normalized = normalize_structure(parsed, slides_count, topic, research)
        return normalized
        
    except Exception as e:
        print(f"Gemini generation failed: {str(e)}")
        # You may want to call generate_fallback_content here
        raise
    # """Generate presentation content using Gemini, grounded with research when available.
    # Enforce: Introduction, Advantages, Disadvantages, Conclusion. If more slides requested, add informative sections.
    # """
    # try:
        
    #     research_text = ''
    #     references = []
    #     if research and research.get('summary'):
    #         research_text = research['summary'][:6000]
    #         if research.get('url'):
    #             references.append(research['url'])

    #     prompt = f"""
    #     You are a professional presentation creator. Build a factual, well-structured PowerPoint outline
    #     about the topic "{topic}".

    #     Ground your content in the following research (if present). Do not fabricate facts or statistics.
    #     Research:
    #     ---
    #     {research_text}
    #     ---

    #     Requirements:
    #     - Produce EXACTLY {slides_count} content slides (title slide is added separately by the app).
    #     - Enforce order and presence:
    #       1) Introduction (3-5 concise sentences tailored to the topic, explain what it is and why it matters)
    #       2) Advantages (4-6 specific, fact-grounded benefits with clear explanations)
    #       3) Disadvantages (4-6 specific, fact-grounded limitations or challenges with clear explanations)
    #       4) Conclusion (3-4 concise sentences summarizing key takeaways and next steps)
    #     - If more than 4 slides requested, add additional informative sections BEFORE Conclusion with clear titles like "Key Concepts", "Use Cases", "Examples", "Statistics", "Trends", "Current Status", "Future Outlook", "Implementation", "Best Practices", or "Challenges". Ensure they are relevant to the topic and grounded in the research.
    #     - Each non-intro/conclusion slide should have 4-6 concise, informative bullets.
    #     - Make content specific to the topic, not generic. Use research data when available.
    #     - Include optional speaker_notes per slide (1-3 sentences) to guide the presenter.
    #     - Include a references array with source URLs at the end (use research URL if available).

    #     Return ONLY valid JSON in this schema:
    #     {{
    #       "title": "string",
    #       "slides": [
    #         {{
    #           "title": "string",
    #           "content": ["bullet 1", "bullet 2", "bullet 3"],
    #           "speaker_notes": "optional short paragraph"
    #         }}
    #       ],
    #       "references": ["https://..."]
    #     }}
    #     """

    #     print(f"Generating content with Gemini for topic: {topic}, slides: {slides_count}")
        
    #     response = model.generate_content(
    #         prompt,
    #         generation_config=genai.types.GenerationConfig(
    #             temperature=0.4,
    #             max_output_tokens=1400
    #         )
    #     )
    #     raw = response.text
    #     print(f"Gemini response received, length: {len(raw)}")
        
    #     # Robust JSON extraction in case model wraps with text
    #     parsed = try_parse_json(raw)
    #     if parsed is None:
    #         print("Failed to parse Gemini response as JSON, using fallback")
    #         raise ValueError("Model returned non-JSON")

    #     # Ensure references include research url if missing
    #     if references:
    #         parsed.setdefault('references', [])
    #         if research['url'] not in parsed['references']:
    #             parsed['references'].append(research['url'])

    #     print(f"Successfully generated {len(parsed.get('slides', []))} slides with Gemini")
        
    #     # Normalize structure strictly to requested format and count
    #     normalized = normalize_structure(parsed, slides_count, topic, research)
    #     return normalized
        
    # except Exception as e:
    #     print(f"Gemini generation failed: {str(e)}")
    #     raise

def try_parse_json(text):
    try:
        return json.loads(text)
    except Exception:
        try:
            start = text.find('{')
            end = text.rfind('}')
            if start != -1 and end != -1 and end > start:
                return json.loads(text[start:end+1])
        except Exception:
            return None
    return None

def fetch_wikipedia_research(topic, language='en'):
    """Fetch a concise summary and url from Wikipedia without API keys."""
    try:
        # Search for best matching page
        search_url = f"https://{language}.wikipedia.org/w/api.php"
        params = {
            'action': 'opensearch',
            'search': topic,
            'limit': 1,
            'namespace': 0,
            'format': 'json'
        }
        r = requests.get(search_url, params=params, timeout=10)
        r.raise_for_status()
        data = r.json()
        if len(data) >= 4 and data[1]:
            page_title = data[1][0]
        else:
            page_title = topic

        # Get summary
        summary_url = f"https://{language}.wikipedia.org/api/rest_v1/page/summary/{page_title}"
        s = requests.get(summary_url, timeout=10)
        if s.status_code == 200:
            sdata = s.json()
            return {
                'summary': sdata.get('extract', ''),
                'url': (sdata.get('content_urls', {}) or {}).get('desktop', {}).get('page')
            }
    except Exception:
        pass
    return {'summary': None, 'url': None}

def normalize_structure(parsed, desired_count, topic, research):
    """Ensure slides include Intro, Advantages, Disadvantages, Conclusion in the correct order,
    and total count equals desired_count. Fill from research when missing.
    """
    slides = parsed.get('slides') or []

    def find_first_with_keyword(keyword_list):
        for idx, s in enumerate(slides):
            title = (s.get('title') or '').lower()
            if any(k in title for k in keyword_list):
                return idx
        return None

    def make_bullets_from_text(text, max_items=6):
        if not text:
            return []
        parts = [p.strip() for p in text.replace('\n', ' ').split('.') if p.strip()]
        return parts[:max_items]

    # Extract or create required sections
    intro_idx = find_first_with_keyword(['intro'])
    adv_idx = find_first_with_keyword(['advantage', 'benefit', 'pros'])
    disadv_idx = find_first_with_keyword(['disadvantage', 'risk', 'cons', 'limitation'])
    concl_idx = find_first_with_keyword(['conclusion', 'closing', 'summary'])

    research_summary = (research or {}).get('summary') or ''

    def pop_slide(idx):
        return slides.pop(idx) if idx is not None and 0 <= idx < len(slides) else None

    intro = pop_slide(intro_idx) or {
        'title': 'Introduction',
        'content': make_bullets_from_text(research_summary, max_items=4) or [f'Overview of {topic}', 'Context and scope', 'Objectives'],
        'speaker_notes': None
    }
    # Adjust index if intro was before advantages
    if adv_idx is not None and intro_idx is not None and adv_idx > intro_idx:
        adv_idx -= 1
    advantages = pop_slide(adv_idx) or {
        'title': 'Advantages',
        'content': make_bullets_from_text(research_summary, max_items=6)[:4] or ['Benefit 1', 'Benefit 2', 'Benefit 3', 'Benefit 4'],
        'speaker_notes': None
    }
    # Recompute disadv index since list may have changed
    disadv_idx = find_first_with_keyword(['disadvantage', 'risk', 'cons', 'limitation'])
    disadvantages = pop_slide(disadv_idx) or {
        'title': 'Disadvantages',
        'content': make_bullets_from_text(research_summary, max_items=6)[4:8] or ['Limitation 1', 'Limitation 2', 'Limitation 3'],
        'speaker_notes': None
    }
    concl_idx = find_first_with_keyword(['conclusion', 'closing', 'summary'])
    conclusion = pop_slide(concl_idx) or {
        'title': 'Conclusion',
        'content': [f'Key takeaways about {topic}', 'Implications', 'Next steps'],
        'speaker_notes': None
    }

    # Remaining slides are "additional"; trim bullets length
    additional = []
    for s in slides:
        if not s or not isinstance(s, dict):
            continue
        title_lower = (s.get('title') or '').lower()
        if any(k in title_lower for k in ['intro', 'advantage', 'benefit', 'pros', 'disadvantage', 'risk', 'cons', 'limitation', 'conclusion', 'closing', 'summary']):
            continue
        bullets = s.get('content') or []
        if len(bullets) < 3:
            enrich = make_bullets_from_text(research_summary)
            bullets = (bullets + enrich)[:6]
        s['content'] = bullets[:6]
        additional.append(s)

    # Build in order: intro, advantages, disadvantages, additional..., conclusion
    ordered = [intro, advantages, disadvantages]
    
    # Calculate how many additional slides we need to reach desired_count
    needed_additional = desired_count - 4  # 4 = intro + advantages + disadvantages + conclusion
    
    # Add existing additional slides first
    while len(ordered) < 3 + needed_additional and additional:
        ordered.append(additional.pop(0))
    
    # If we still need more slides, synthesize them from research
    synth_titles = ['Key Concepts', 'Use Cases', 'Examples', 'Statistics', 'Trends', 'Current Status', 'Future Outlook', 'Implementation', 'Best Practices', 'Challenges']
    synth_idx = 0
    while len(ordered) < 3 + needed_additional and synth_idx < len(synth_titles):
        bullets = make_bullets_from_text(research_summary) or [f'Point {i}' for i in range(1, 5)]
        ordered.append({'title': synth_titles[synth_idx], 'content': bullets[:5], 'speaker_notes': None})
        synth_idx += 1
    
    # Append conclusion at the end
    ordered.append(conclusion)
    
    # Ensure we have exactly the requested number of slides
    if len(ordered) < desired_count:
        # Fill remaining with generic informative slides
        remaining = desired_count - len(ordered)
        for i in range(remaining):
            bullets = make_bullets_from_text(research_summary) or [f'Key point {i+1}', f'Important aspect {i+1}', f'Notable feature {i+1}']
            ordered.append({
                'title': f'Additional Information {i+1}',
                'content': bullets[:4],
                'speaker_notes': None
            })
    elif len(ordered) > desired_count:
        # Trim to exact count, keeping intro, advantages, disadvantages, conclusion
        ordered = ordered[:desired_count-1] + [conclusion]

    return {
        'title': parsed.get('title') or topic,
        'slides': ordered,
        'references': parsed.get('references') or (([research.get('url')] if research and research.get('url') else []))
    }

    
    # Find the best matching topic or use generic professional content
    topic_lower = topic.lower()
    matched_content = None
    
    for key, content in professional_content.items():
        if key in topic_lower or any(word in topic_lower for word in key.split()):
            matched_content = content
            break
    
    if not matched_content:
        # Generate generic professional content based on topic
        matched_content = {
            "intro": [
                f"Welcome to our comprehensive analysis of {topic}",
                f"This presentation explores the fundamental concepts and applications of {topic}",
                f"We will examine the key aspects, benefits, challenges, and future implications",
                f"Understanding {topic} is crucial for success in today's competitive landscape"
            ],
            "advantages": [
                f"Enhanced efficiency and productivity in {topic}-related applications",
                f"Improved accuracy and precision through systematic approaches",
                f"Cost-effective solutions for complex {topic} challenges",
                f"Scalability and adaptability to different business scenarios",
                f"Innovation and competitive advantage in {topic} implementation"
            ],
            "disadvantages": [
                f"Initial implementation costs and resource requirements for {topic}",
                f"Complexity in understanding and applying {topic} principles",
                f"Potential risks and limitations in certain {topic} contexts",
                f"Need for continuous updates and maintenance of {topic} systems",
                f"Dependency on external factors and market conditions"
            ],
            "additional": [
                f"Key Concepts: Fundamental principles and theories of {topic}",
                f"Use Cases: Real-world applications and implementation examples",
                f"Best Practices: Industry standards and proven methodologies",
                f"Trends: Current developments and future outlook in {topic}",
                f"Implementation: Strategic planning and execution frameworks"
            ]
        }
    
    # Base structure with professional content
    base_slides = [
        {
            "title": "Introduction",
            "content": matched_content["intro"],
            "speaker_notes": f"Set the stage for understanding {topic}. Explain its significance and relevance in today's context. Highlight key learning objectives."
        },
        {
            "title": "Advantages & Benefits",
            "content": matched_content["advantages"],
            "speaker_notes": f"Present the compelling benefits of {topic}. Use specific examples and data points to support each advantage. Address ROI and business impact."
        },
        {
            "title": "Challenges & Considerations",
            "content": matched_content["disadvantages"],
            "speaker_notes": f"Discuss the realistic challenges and limitations of {topic}. Provide mitigation strategies and risk management approaches."
        }
    ]
    
    # Additional professional slides if more requested
    additional_titles = [
        "Key Concepts & Fundamentals",
        "Implementation Strategies", 
        "Best Practices & Standards",
        "Case Studies & Examples",
        "Current Trends & Developments",
        "Future Outlook & Predictions",
        "Risk Management & Mitigation",
        "Performance Metrics & KPIs",
        "Technology & Tools",
        "Industry Applications"
    ]
    
    additional_slides = []
    for i in range(max(0, slides_count - 4)):
        title = additional_titles[i % len(additional_titles)]
        
        # Always generate professional content for additional slides
        if title == "Key Concepts & Fundamentals":
            content = [
                f"Core principles and foundational theories of {topic}",
                f"Essential frameworks and methodologies for understanding",
                f"Key terminology and industry-specific language",
                f"Fundamental concepts that drive success in {topic}"
            ]
        elif title == "Implementation Strategies":
            content = [
                f"Step-by-step approach to implementing {topic} solutions",
                f"Resource allocation and timeline planning considerations",
                f"Change management and stakeholder engagement strategies",
                f"Risk assessment and mitigation planning approaches"
            ]
        elif title == "Best Practices & Standards":
            content = [
                f"Industry-recognized standards and quality benchmarks",
                f"Proven methodologies and success patterns",
                f"Common pitfalls and how to avoid them",
                f"Continuous improvement and optimization strategies"
            ]
        elif title == "Case Studies & Examples":
            content = [
                f"Real-world examples of successful {topic} implementation",
                f"Lessons learned from industry leaders and innovators",
                f"Quantifiable results and performance improvements",
                f"Adaptable strategies for different organizational contexts"
            ]
        elif title == "Current Trends & Developments":
            content = [
                f"Emerging technologies and innovative approaches in {topic}",
                f"Market dynamics and competitive landscape changes",
                f"Regulatory updates and compliance requirements",
                f"Industry disruption and transformation opportunities"
            ]
        elif title == "Future Outlook & Predictions":
            content = [
                f"Projected growth and market expansion in {topic}",
                f"Technology evolution and capability enhancements",
                f"Strategic opportunities and competitive advantages",
                f"Long-term planning and investment considerations"
            ]
        elif title == "Risk Management & Mitigation":
            content = [
                f"Identification of potential risks in {topic} implementation",
                f"Assessment of impact and probability for each risk",
                f"Development of comprehensive mitigation strategies",
                f"Monitoring and response protocols for risk events"
            ]
        elif title == "Performance Metrics & KPIs":
            content = [
                f"Key performance indicators for measuring {topic} success",
                f"Data collection and analysis methodologies",
                f"Benchmarking against industry standards and competitors",
                f"Reporting frameworks and stakeholder communication"
            ]
        elif title == "Technology & Tools":
            content = [
                f"Essential technologies and platforms for {topic} success",
                f"Integration requirements and compatibility considerations",
                f"Training and skill development requirements",
                f"Technology roadmap and upgrade planning"
            ]
        else:  # Industry Applications
            content = [
                f"Cross-industry applications and use cases for {topic}",
                f"Sector-specific considerations and adaptations",
                f"Scalability and customization requirements",
                f"Partnership and collaboration opportunities"
            ]
        
        additional_slides.append({
            "title": title,
            "content": content,
            "speaker_notes": f"Deep dive into {title.lower()}. Provide actionable insights and practical examples that attendees can apply immediately."
        })
    
    # Professional conclusion
    conclusion = {
        "title": "Conclusion & Next Steps",
        "content": [
            f"Key takeaways and strategic insights about {topic}",
            f"Critical success factors for successful {topic} implementation",
            f"Immediate action items and recommended next steps",
            f"Resources and references for continued learning and development"
        ],
        "speaker_notes": f"Summarize the key points and provide a clear roadmap for implementation. Encourage questions and discussion. Thank participants for their engagement."
    }
    
    # Combine all slides
    all_slides = base_slides + additional_slides + [conclusion]
    
    # Ensure exact count
    all_slides = all_slides[:slides_count]
    
    print(f"Professional fallback content generated: {len(all_slides)} slides")
    
    return {
        "title": f"Professional Guide: {topic} - Strategies, Implementation & Best Practices",
        "slides": all_slides,
        "references": ([research.get('url')] if research and research.get('url') else [])
    }

def create_presentation(topic, content, theme_file):
    """Create PowerPoint presentation from content, with speaker notes and references."""
    prs = Presentation(rf"themes/{theme_file}")
    
    # Set slide dimensions
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Remove all slides from the loaded theme
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = content.get('title', f"Presentation on {topic}")
    subtitle.text = f"Generated by AI Presentation Creator"
    # Style the title slide
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)
    
    # Content slides
    for slide_data in content.get('slides', []):
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title
        title = slide.shapes.title
        title.text = slide_data.get('title', 'Slide Title')
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)
        
        # Add content
        content_placeholder = slide.placeholders[1]
        text_frame = content_placeholder.text_frame
        
        for i, point in enumerate(slide_data.get('content', [])):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = f"• {point}"
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(52, 73, 94)
            p.space_after = Pt(12)

        # Optional speaker notes
        speaker_notes = slide_data.get('speaker_notes')
        if speaker_notes:
            notes_tf = slide.notes_slide.notes_text_frame
            notes_tf.text = speaker_notes

    # References slide
    refs = content.get('references') or []
    if refs:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "References"
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)
        content_placeholder = slide.placeholders[1]
        text_frame = content_placeholder.text_frame
        for idx, ref in enumerate(refs):
            p = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()
            p.text = ref
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(52, 73, 94)
            p.space_after = Pt(8)
    
    # Save to temporary file
    temp_dir = tempfile.gettempdir()
    pptx_path = os.path.join(temp_dir, f"{topic.replace(' ', '_')}_presentation.pptx")
    prs.save(pptx_path)
    
    return pptx_path

@app.route('/api/health', methods=['GET'])
def health_check():
    return jsonify({'status': 'healthy', 'message': 'PPT Generator API is running'})

if __name__ == '__main__':
    run("app:asgi_app", port=8000, reload=True)